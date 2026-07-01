"""PowerPoint presentation generation for the analytical workflow."""

from __future__ import annotations

from pathlib import Path

from PIL import Image

EXCLUDED_FIGURE_NAMES = {
    "Figure_2B_ChangeAgreement_2009_2013.png",
    "Figure_2B_ChangeAgreement_2013_2018.png",
    "Figure_2B_ChangeAgreement_2018_2024.png",
    "Figure_04_CTrees_Snapshot_x_MB_LULC.png",
    "Figure_02_AreaBar_30m.png",
    "Figure_08_ForestChangeArea_TimeSeries_30m.png",
    "Figure_11_TemporalReversal_Area_30m.png",
}


def _add_picture_fit(slide, image_path: Path, left, top, max_width, max_height):
    """Add an image to a slide without distorting its aspect ratio."""
    with Image.open(image_path) as image:
        pixel_width, pixel_height = image.size
    if pixel_width <= 0 or pixel_height <= 0:
        return slide.shapes.add_picture(str(image_path), left, top, width=max_width)

    image_ratio = pixel_width / pixel_height
    box_ratio = max_width / max_height
    if image_ratio >= box_ratio:
        width = max_width
        height = int(max_width / image_ratio)
    else:
        height = max_height
        width = int(max_height * image_ratio)

    x = left + int((max_width - width) / 2)
    y = top + int((max_height - height) / 2)
    return slide.shapes.add_picture(str(image_path), x, y, width=width, height=height)


def _patch_pptx_picture_aspect_fit() -> None:
    """Preserve image aspect ratio for slide pictures inserted into fixed boxes."""
    try:
        from pptx.shapes.shapetree import SlideShapes
    except Exception:
        return
    if getattr(SlideShapes.add_picture, "_mapbiomas_ctrees_aspect_fit", False):
        return

    original_add_picture = SlideShapes.add_picture

    def add_picture_aspect_fit(self, image_file, left, top, width=None, height=None):
        if width is None or height is None:
            return original_add_picture(self, image_file, left, top, width=width, height=height)
        try:
            with Image.open(image_file) as image:
                pixel_width, pixel_height = image.size
            if pixel_width <= 0 or pixel_height <= 0:
                return original_add_picture(self, image_file, left, top, width=width, height=height)
            image_ratio = pixel_width / pixel_height
            box_ratio = width / height
            if image_ratio >= box_ratio:
                fitted_width = width
                fitted_height = int(width / image_ratio)
            else:
                fitted_height = height
                fitted_width = int(height * image_ratio)
            fitted_left = left + int((width - fitted_width) / 2)
            fitted_top = top + int((height - fitted_height) / 2)
            return original_add_picture(
                self,
                image_file,
                fitted_left,
                fitted_top,
                width=fitted_width,
                height=fitted_height,
            )
        except Exception:
            return original_add_picture(self, image_file, left, top, width=width, height=height)

    add_picture_aspect_fit._mapbiomas_ctrees_aspect_fit = True
    SlideShapes.add_picture = add_picture_aspect_fit


_patch_pptx_picture_aspect_fit()


def build_powerpoint_presentation(
    table_directory: Path,
    figure_directory: Path,
    output_path: Path,
) -> Path:
    """Create one presentation slide per generated figure."""
    try:
        from pptx import Presentation
        from pptx.enum.text import PP_ALIGN as pptx_align
        from pptx.util import Inches as pptx_inches
        from pptx.util import Pt as pptx_pt
    except ModuleNotFoundError as exc:
        raise RuntimeError(
            "PowerPoint generation requires python-pptx. Install dependencies with "
            "`python -m pip install -r requirements.txt`."
        ) from exc
    global Inches, PP_ALIGN, Pt
    Inches = pptx_inches
    PP_ALIGN = pptx_align
    Pt = pptx_pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    _add_title_slide(presentation)

    for figure_path in sorted(figure_directory.glob("*.png")):
        _add_figure_slide(presentation, figure_path)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return output_path


def _add_title_slide(presentation) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(slide, "MapBiomas vs. CTrees Cross-Validation Technical Presentation")
    _add_textbox(
        slide,
        "Pará BVP2 | TerraCarbon / Verra VMD0055 v1.1 | VT0007 v1.0",
        Inches(0.8),
        Inches(2.0),
        Inches(11.8),
        Inches(0.7),
        24,
    )
    _add_footer(slide, "Source: VMD0055 v1.1 and VT0007 v1.0 analytical workflow.")


def _add_figure_slide(presentation, figure_path: Path) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _add_title(slide, _clean_title(figure_path.stem))
    slide.shapes.add_picture(str(figure_path), Inches(0.7), Inches(1.05), width=Inches(11.95), height=Inches(5.65))
    _add_footer(slide, _verra_reference_for_name(figure_path.stem))


def _add_title(slide, title: str) -> None:
    _add_textbox(slide, title, Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.55), 24, bold=True)


def _add_textbox(slide, text: str, x, y, width, height, font_size: int, bold: bool = False) -> None:
    shape = slide.shapes.add_textbox(x, y, width, height)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.alignment = PP_ALIGN.LEFT


def _add_footer(slide, text: str) -> None:
    _add_textbox(slide, text, Inches(0.45), Inches(7.05), Inches(12.45), Inches(0.25), 8)


def _clean_title(value: str) -> str:
    return " ".join(value.replace("_", " ").replace("-", " to ").split()).title()


def _verra_reference_for_name(name: str) -> str:
    lower = name.lower()
    if "fcbm" in lower or "udefa" in lower:
        return "Source: VT0007 v1.0, Table 1 and Section Data Requirements; VMD0055 v1.1, Tables 15 and 16."
    if "change" in lower or "agreement" in lower:
        return "Source: VMD0055 v1.1, Table 15 for interpreted CTrees classes."
    return "Source: MapBiomas Collection 10.1 and CTrees reference products."
